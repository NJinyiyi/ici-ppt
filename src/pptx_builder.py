from __future__ import annotations

from pathlib import Path


def build_pptx_from_pngs(png_paths: list[Path], output_path: Path, title: str) -> Path:
    """Build a stable 16:9 PPTX by inserting each rendered PNG as a full-slide image.

    The previous implementation wrote OOXML by hand. PowerPoint is strict about
    package relationships and slide metadata, so we use python-pptx here to avoid
    repair prompts when opening the generated deck.
    """
    if not png_paths:
        raise ValueError("No PNG paths supplied for PPTX generation.")

    from pptx import Presentation
    from pptx.util import Inches

    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = title
    prs.core_properties.author = "ici-ppt"
    prs.core_properties.last_modified_by = "ici-ppt"

    blank_layout = prs.slide_layouts[6]
    for png in png_paths:
        if not png.exists():
            raise FileNotFoundError(f"Missing slide PNG: {png}")
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(png),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    # Remove the default empty slide if the template ever creates one. The stock
    # python-pptx Presentation starts with zero slides, so this is normally a no-op.
    prs.save(output_path)
    return output_path
