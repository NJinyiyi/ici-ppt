from __future__ import annotations

from pathlib import Path
from typing import Iterable

from planner import Slide


FONT_LIGHT = "Alibaba PuHuiTi Light"
FONT_REGULAR = "Alibaba PuHuiTi"
FONT_MEDIUM = "Alibaba PuHuiTi Medium"
FONT_BOLD = "Alibaba PuHuiTi Bold"
FONT_HEAVY = "Alibaba PuHuiTi Heavy"

DEEP_BLUE = "1B1464"
PURPLE_BLUE = "25148A"
MAIN_BLUE = "006FD6"
CYAN_GREEN = "66E6C3"
DARK_TEXT = "111827"
SECONDARY_TEXT = "4B5563"
LIGHT_BG = "F8FAFC"
WHITE = "FFFFFF"
SOFT_LINE = "DDE7F3"


def build_editable_pptx(slides: list[Slide], output_path: Path, title: str) -> Path:
    """Build a stable editable deck using python-pptx native shapes and text."""
    if not slides:
        raise ValueError("No slides supplied for editable PPTX generation.")

    from pptx import Presentation
    from pptx.util import Inches

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = title
    prs.core_properties.author = "ici-ppt"
    prs.core_properties.last_modified_by = "ici-ppt"

    blank = prs.slide_layouts[6]
    for slide_model in slides:
        slide = prs.slides.add_slide(blank)
        draw_slide(slide, slide_model)

    prs.save(output_path)
    return output_path


def draw_slide(slide, model: Slide) -> None:
    layout = model.layout
    if layout in {"cover", "toc", "section", "closing"}:
        draw_gradient_background(slide)
        add_text(slide, "ICI Lab · Zhejiang University", 120, 72, 760, 38, 18, WHITE, FONT_BOLD)
        add_tri_squares(slide, gradient=True)
    else:
        add_rect(slide, 0, 0, 1920, 1080, LIGHT_BG, no_line=True)
        add_text(slide, "ICI Lab", 120, 72, 260, 34, 18, DEEP_BLUE, FONT_BOLD)
        add_tri_squares(slide, gradient=False)

    if layout == "cover":
        add_text(slide, model.title, 120, 286, 1280, 240, 52, WHITE, FONT_HEAVY)
        add_text(slide, model.data.get("subtitle", ""), 120, 548, 1210, 88, 25, "E6F5FF", FONT_REGULAR)
        add_text(slide, model.data.get("meta", ""), 120, 930, 980, 46, 19, "DCEFFF", FONT_REGULAR)
    elif layout == "toc":
        add_text(slide, "Contents", 120, 160, 640, 82, 44, WHITE, FONT_HEAVY)
        draw_toc(slide, model.data.get("items", []))
    elif layout == "section":
        add_text(slide, model.data.get("section_no", "01"), 120, 230, 430, 145, 82, WHITE, FONT_HEAVY)
        add_text(slide, wrap_text(model.title, 18, 2), 120, 405, 1280, 140, 38, WHITE, FONT_HEAVY, line_spacing=0.95)
        add_text(slide, wrap_text(model.data.get("subtitle", ""), 36, 2), 120, 592, 1220, 82, 21, "E1F5FF", FONT_REGULAR, line_spacing=1.05)
    elif layout == "closing":
        add_text(slide, model.title, 120, 245, 1370, 180, 50, WHITE, FONT_HEAVY)
        add_text(slide, model.data.get("message", ""), 120, 470, 1180, 180, 26, "E6F5FF", FONT_REGULAR)
        add_text(slide, model.data.get("contact", ""), 120, 900, 980, 48, 20, "E1F5FF", FONT_REGULAR)
    else:
        draw_content_header(slide, model)
        draw_body(slide, model)
        add_text(slide, model.data.get("footer", "ICI Lab · Zhejiang University"), 120, 998, 700, 32, 16, "64748B", FONT_REGULAR)


def draw_content_header(slide, model: Slide) -> None:
    add_text(slide, model.data.get("kicker", "ICI Lab Report"), 120, 66, 700, 34, 14, MAIN_BLUE, FONT_BOLD)
    add_text(slide, model.title, 120, 124, 1260, 88, 30, DEEP_BLUE, FONT_HEAVY)
    add_rect(slide, 120, 224, 150, 6, MAIN_BLUE, no_line=True)


def draw_body(slide, model: Slide) -> None:
    if model.layout == "two_column":
        add_card(slide, 120, 310, 790, 540, model.data.get("left_title", "Challenge"), model.data.get("left_items", []))
        add_card(slide, 1010, 310, 790, 540, model.data.get("right_title", "Response"), model.data.get("right_items", []), accent=CYAN_GREEN)
    elif model.layout == "process":
        steps = model.data.get("steps", [])[:4]
        count = max(1, len(steps))
        gap = 42 if count <= 3 else 26
        card_w = int((1680 - gap * (count - 1)) / count)
        for idx, step in enumerate(steps, start=1):
            add_process_card(slide, 120 + (idx - 1) * (card_w + gap), 330, card_w, 410, f"{idx:02d}", step)
        add_text(slide, wrap_text(model.data.get("note", ""), 46, 3), 120, 805, 1180, 105, 19, SECONDARY_TEXT, FONT_REGULAR, line_spacing=1.05)
    elif model.layout == "image":
        add_figure_placeholder(slide, 120, 310, 1040, 585, model.data.get("figure_title", "Figure / Prototype / Result"), model.data.get("figure_note", ""))
        add_bullets(slide, model.data.get("items", []), 1260, 330, 520, 430, 19)
    elif model.layout == "summary":
        for idx, card in enumerate(model.data.get("cards", [])[:3]):
            card_title, body = card
            add_card(slide, 120 + idx * 570, 318, 530, 418, card_title, [body], accent=CYAN_GREEN)
    else:
        add_bullets(slide, model.data.get("items", []), 120, 300, 1180, 500, 21)
        add_card(slide, 1420, 300, 380, 400, "Key Point", [model.data.get("highlight", model.title)], accent=MAIN_BLUE)


def draw_toc(slide, items: Iterable[str]) -> None:
    toc_items = list(items)[:4]
    rows = 2 if len(toc_items) <= 4 else 3
    for idx, item in enumerate(toc_items, start=1):
        col = 0 if idx <= rows else 1
        row = idx - 1 if idx <= rows else idx - rows - 1
        x = 120 + col * 860
        y = 330 + row * 150
        add_text(slide, f"{idx:02d}", x, y, 86, 52, 27, WHITE, FONT_HEAVY)
        add_text(slide, wrap_text(str(item), 16, 2), x + 112, y + 2, 610, 74, 22, "EBFAFF", FONT_BOLD, line_spacing=0.98)
        add_rect(slide, x + 112, y + 92, 240, 3, "EBFAFF", no_line=True)


def draw_gradient_background(slide) -> None:
    add_gradient_rect(
        slide,
        0,
        0,
        1920,
        1080,
        stops=[(0, PURPLE_BLUE), (48000, MAIN_BLUE), (100000, CYAN_GREEN)],
        angle=135,
    )


def add_tri_squares(slide, gradient: bool) -> None:
    colors = [WHITE, "E6F5FF", "CDEBEB"] if gradient else [PURPLE_BLUE, MAIN_BLUE, CYAN_GREEN]
    for idx, color in enumerate(colors):
        add_rect(slide, 1710 + idx * 40, 988, 28, 28, color, no_line=True)


def add_card(slide, x: int, y: int, w: int, h: int, title: str, items: Iterable[str], accent: str | None = None) -> None:
    add_rect(slide, x, y, w, h, WHITE, line=SOFT_LINE)
    if accent:
        add_rect(slide, x, y, w, 8, accent, no_line=True)
    add_text(slide, wrap_text(str(title), max(8, int((w - 68) / 34)), 2), x + 34, y + 34, w - 68, 70, 24, DEEP_BLUE, FONT_HEAVY)
    add_bullets(slide, list(items)[:4], x + 34, y + 132, w - 68, h - 160, 18)


def add_process_card(slide, x: int, y: int, w: int, h: int, number: str, step: str) -> None:
    add_rect(slide, x, y, w, h, WHITE, line=SOFT_LINE)
    add_rect(slide, x, y, w, 8, MAIN_BLUE, no_line=True)
    add_text(slide, number, x + 34, y + 38, w - 68, 58, 28, DEEP_BLUE, FONT_HEAVY)
    chars = max(8, int((w - 68) / 34))
    add_text(slide, wrap_text(step, chars, 7), x + 34, y + 128, w - 68, h - 158, 18, DARK_TEXT, FONT_MEDIUM, line_spacing=1.02)


def add_figure_placeholder(slide, x: int, y: int, w: int, h: int, title: str, note: str) -> None:
    add_rect(slide, x, y, w, h, "F1F7FF", line=SOFT_LINE)
    add_rect(slide, x + 28, y + 28, w - 56, h - 102, "FFFFFF", line=SOFT_LINE)
    add_text(slide, wrap_text(title, 22, 2), x + 120, y + 235, w - 240, 78, 25, DEEP_BLUE, FONT_HEAVY, align="center", line_spacing=0.95)
    add_text(slide, wrap_text(note, 36, 2), x + 160, y + 318, w - 320, 58, 16, SECONDARY_TEXT, FONT_REGULAR, align="center", line_spacing=1.05)
    add_rect(slide, x + 28, y + h - 52, w - 56, 1, SOFT_LINE, no_line=True)
    add_text(slide, "Figure placeholder", x + 42, y + h - 42, w - 84, 28, 13, "64748B", FONT_REGULAR)


def add_bullets(slide, items: Iterable[str], x: int, y: int, w: int, h: int, size: int) -> None:
    item_list = list(items)
    chars = max(8, int(w / (size * 1.75)))
    wrapped_items = []
    max_items = min(5, len(item_list))
    max_lines = max(1, int(h / (size * 2.0)) // max(1, max_items))
    for item in item_list[:5]:
        lines = wrap_text(str(item), chars, max_lines).splitlines()
        if not lines:
            continue
        wrapped_items.append("• " + lines[0])
        wrapped_items.extend("  " + line for line in lines[1:])
    text = "\n".join(wrapped_items)
    add_text(slide, text, x, y, w, h, size, DARK_TEXT, FONT_REGULAR, line_spacing=1.18)


def wrap_text(text: str, max_chars: int, max_lines: int) -> str:
    text = " ".join(str(text).split())
    if not text:
        return ""
    lines: list[str] = []
    current = ""
    current_width = 0.0
    for char in text:
        width = 0.55 if char.isascii() and char not in "，。；：、！？（）" else 1.0
        if current and current_width + width > max_chars:
            lines.append(current.rstrip())
            current = char
            current_width = width
            if len(lines) >= max_lines:
                break
        else:
            current += char
            current_width += width
    if len(lines) < max_lines and current:
        lines.append(current.rstrip())
    if len(lines) > max_lines:
        lines = lines[:max_lines]
    if len(lines) == max_lines and len("".join(lines)) < len(text):
        lines[-1] = lines[-1].rstrip("，。；：、,. ") + "…"
    return "\n".join(lines)


def add_rect(slide, x: int, y: int, w: int, h: int, fill: str, line: str | None = None, no_line: bool = False):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x), px(y), px(w), px(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill)
    if no_line:
        shape.line.fill.background()
    elif line:
        shape.line.color.rgb = RGBColor.from_string(line)
        shape.line.width = px(1)
    else:
        shape.line.fill.background()
    return shape


def add_gradient_rect(slide, x: int, y: int, w: int, h: int, stops: list[tuple[int, str]], angle: int = 135):
    shape = add_rect(slide, x, y, w, h, stops[0][1], no_line=True)
    apply_gradient_fill(shape, stops, angle)
    return shape


def apply_gradient_fill(shape, stops: list[tuple[int, str]], angle: int) -> None:
    from pptx.oxml.xmlchemy import OxmlElement
    from pptx.oxml.ns import qn

    sp_pr = shape._element.spPr
    for fill_tag in ("a:solidFill", "a:gradFill", "a:noFill", "a:pattFill", "a:blipFill"):
        node = sp_pr.find(qn(fill_tag))
        if node is not None:
            sp_pr.remove(node)

    grad_fill = OxmlElement("a:gradFill")
    grad_fill.set("rotWithShape", "1")

    gs_lst = OxmlElement("a:gsLst")
    for position, color in stops:
        gs = OxmlElement("a:gs")
        gs.set("pos", str(position))
        srgb = OxmlElement("a:srgbClr")
        srgb.set("val", color)
        gs.append(srgb)
        gs_lst.append(gs)
    grad_fill.append(gs_lst)

    lin = OxmlElement("a:lin")
    lin.set("ang", str(angle * 60000))
    lin.set("scaled", "1")
    grad_fill.append(lin)

    children = list(sp_pr)
    ln = sp_pr.find(qn("a:ln"))
    if ln is not None:
        sp_pr.insert(children.index(ln), grad_fill)
    else:
        sp_pr.append(grad_fill)


def add_text(
    slide,
    text: str,
    x: int,
    y: int,
    w: int,
    h: int,
    size: int,
    color: str,
    font_name: str,
    align: str = "left",
    line_spacing: float | None = None,
):
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
    from pptx.util import Pt

    box = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)

    lines = str(text).splitlines() or [""]
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = ""
        p.alignment = PP_ALIGN.CENTER if align == "center" else PP_ALIGN.LEFT
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.color.rgb = RGBColor.from_string(color)
        if font_name in {FONT_BOLD, FONT_HEAVY}:
            run.font.bold = True
        set_east_asian_font(run, font_name)
    return box


def set_east_asian_font(run, font_name: str) -> None:
    from pptx.oxml.xmlchemy import OxmlElement
    from pptx.oxml.ns import qn

    r_pr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        node = r_pr.find(qn(tag))
        if node is None:
            node = OxmlElement(tag)
            r_pr.append(node)
        node.set("typeface", font_name)


def px(value: float):
    from pptx.util import Inches

    return Inches(value / 144)
