from __future__ import annotations

from pathlib import Path
from urllib.parse import unquote, urlparse
from typing import Any

from editable_pptx_builder import (
    CYAN_GREEN,
    DEEP_BLUE,
    FONT_BOLD,
    FONT_HEAVY,
    FONT_MEDIUM,
    FONT_REGULAR,
    LIGHT_BG,
    MAIN_BLUE,
    PURPLE_BLUE,
    SOFT_LINE,
    WHITE,
    add_gradient_rect,
    px,
    set_east_asian_font,
)


def build_hybrid_pptx(layouts: list[dict[str, Any]], output_path: Path, title: str) -> Path:
    """Rebuild HTML-authored slides as editable PowerPoint shapes.

    The HTML is still the source of visual layout. Playwright extracts each
    marked DOM node's computed box and typography; this builder recreates those
    nodes as native PowerPoint text boxes and shapes.
    """
    if not layouts:
        raise ValueError("No extracted HTML layouts supplied for hybrid PPTX generation.")

    from pptx import Presentation
    from pptx.util import Inches

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = title
    prs.core_properties.author = "ici-ppt"
    prs.core_properties.last_modified_by = "ici-ppt"

    blank = prs.slide_layouts[6]
    for layout in layouts:
        slide = prs.slides.add_slide(blank)
        is_gradient = draw_background(slide, layout)
        for item in layout.get("items", []):
            draw_item(slide, item, is_gradient=is_gradient)

    prs.save(output_path)
    return output_path


def draw_background(slide, layout: dict[str, Any]) -> bool:
    classes = set(layout.get("slide", {}).get("classes", []))
    if "gradient" in classes:
        add_gradient_rect(
            slide,
            0,
            0,
            1920,
            1080,
            stops=[(0, PURPLE_BLUE), (48000, MAIN_BLUE), (100000, CYAN_GREEN)],
            angle=135,
        )
        return True
    else:
        add_rect(slide, 0, 0, 1920, 1080, LIGHT_BG, no_line=True)
        return False


def draw_item(slide, item: dict[str, Any], is_gradient: bool = False) -> None:
    kind = item.get("type")
    role = item.get("role", "")
    if kind in {"text", "bullet"}:
        add_dom_text(slide, item, bullet=kind == "bullet", is_gradient=is_gradient)
    elif kind == "image":
        add_dom_image(slide, item)
    elif kind == "card":
        add_dom_card(slide, item)
    elif kind == "rect":
        style = item.get("style", {})
        fill = style.get("backgroundColor") or role_fill(role)
        add_rect_from_item(slide, item, fill=fill, no_line=True)


def add_dom_card(slide, item: dict[str, Any]) -> None:
    style = item.get("style", {})
    fill = style.get("backgroundColor") or WHITE
    line = style.get("borderColor") or SOFT_LINE
    shape = add_rect_from_item(slide, item, fill=fill, line=line)
    try:
        shape.shadow.inherit = False
    except Exception:
        pass

    border_top = float(style.get("borderTopWidth") or 0)
    if border_top >= 4:
        box = item["box"]
        accent = style.get("borderColor") or MAIN_BLUE
        add_rect(slide, box["x"], box["y"], box["w"], border_top, accent, no_line=True)


def add_dom_text(slide, item: dict[str, Any], bullet: bool = False, is_gradient: bool = False):
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
    from pptx.util import Pt

    box = item["box"]
    style = item.get("style", {})
    text = item.get("text", "")
    role = item.get("role", "")
    if bullet:
        text = "• " + text
        box = {**box, "x": box["x"] - 4, "w": box["w"] + 14}

    tb = slide.shapes.add_textbox(px(box["x"]), px(box["y"]), px(box["w"]), px(box["h"]))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER if style.get("textAlign") == "center" else PP_ALIGN.LEFT
    line_height = style.get("lineHeight")
    if line_height:
        p.line_spacing = float(line_height)

    run = p.add_run()
    run.text = text
    font = font_for(float(style.get("fontWeight") or 400))
    run.font.name = font
    run.font.size = Pt(max(8, min(86, float(style.get("fontSize") or 22))))
    color = style.get("color") or default_text_color(role)
    if is_gradient and is_dark_color(color):
        color = gradient_text_color(role)
    run.font.color.rgb = RGBColor.from_string(color)
    if font in {FONT_BOLD, FONT_HEAVY}:
        run.font.bold = True
    set_east_asian_font(run, font)
    return tb


def add_dom_image(slide, item: dict[str, Any]) -> None:
    src = item.get("src", "")
    path = path_from_file_url(src)
    if not path or not path.exists():
        return
    box = item["box"]
    slide.shapes.add_picture(str(path), px(box["x"]), px(box["y"]), width=px(box["w"]), height=px(box["h"]))


def add_rect_from_item(slide, item: dict[str, Any], fill: str, line: str | None = None, no_line: bool = False):
    box = item["box"]
    return add_rect(slide, box["x"], box["y"], box["w"], box["h"], fill, line=line, no_line=no_line)


def add_rect(slide, x: float, y: float, w: float, h: float, fill: str, line: str | None = None, no_line: bool = False):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x), px(y), px(w), px(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill)
    if no_line:
        shape.line.fill.background()
    elif line:
        shape.line.color.rgb = RGBColor.from_string(line)
        shape.line.width = px(1)
    else:
        shape.line.fill.background()
    return shape


def font_for(weight: float) -> str:
    if weight >= 800:
        return FONT_HEAVY
    if weight >= 700:
        return FONT_BOLD
    if weight >= 600:
        return FONT_MEDIUM
    return FONT_REGULAR


def path_from_file_url(src: str) -> Path | None:
    if not src:
        return None
    parsed = urlparse(src)
    if parsed.scheme == "file":
        return Path(unquote(parsed.path))
    if parsed.scheme == "":
        return Path(src)
    return None


def is_dark_color(color: str | None) -> bool:
    if not color or len(color) != 6:
        return True
    try:
        r = int(color[0:2], 16)
        g = int(color[2:4], 16)
        b = int(color[4:6], 16)
    except ValueError:
        return True
    return (0.2126 * r + 0.7152 * g + 0.0722 * b) < 110


def gradient_text_color(role: str) -> str:
    if role in {"cover-subtitle", "cover-meta", "section-subtitle", "closing-message", "closing-contact"}:
        return "E6F5FF"
    return WHITE


def role_fill(role: str) -> str:
    if role == "rule":
        return MAIN_BLUE
    if role == "tri-square":
        return MAIN_BLUE
    return MAIN_BLUE


def default_text_color(role: str) -> str:
    if role in {"lab-mark", "cover-title", "cover-subtitle", "cover-meta", "toc-heading", "toc-number", "toc-title", "section-number", "section-title", "section-subtitle", "closing-title", "closing-message", "closing-contact"}:
        return WHITE
    if role in {"kicker", "key-card-label", "process-number", "summary-title"}:
        return MAIN_BLUE
    return DEEP_BLUE
